"""
HowTune — Apple Keynote 風プレゼンデッキ
ミニマル · 角丸カード · グラデーション主体 · 透明感
10分 (14スライド) + Q&A対策 6スライド
"""
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from lxml import etree

# ── Palette ──────────────────────────────────────────────────
BLACK  = RGBColor(0x00, 0x00, 0x00)
WHITE  = RGBColor(0xFF, 0xFF, 0xFF)
GRAY   = RGBColor(0x8E, 0x8E, 0x93)
LGRAY  = RGBColor(0xB8, 0xB8, 0xBE)
BLUE   = RGBColor(0x0A, 0x84, 0xFF)
LBLUE  = RGBColor(0x64, 0xD2, 0xFF)
GREEN  = RGBColor(0x30, 0xD1, 0x58)
RED    = RGBColor(0xFF, 0x45, 0x3A)
AMBER  = RGBColor(0xFF, 0x9F, 0x0A)
PURPLE = RGBColor(0xBF, 0x5A, 0xF2)
PINK   = RGBColor(0xFF, 0x37, 0x5F)
GLASS  = RGBColor(0x1C, 0x1C, 0x1E)
DBLUE  = RGBColor(0x06, 0x12, 0x24)

W  = Inches(13.333)
H  = Inches(7.5)
ML = Inches(1.0)
CW = W - ML * 2

import os
OUT = "docs/slides/howtune_final.pptx"

# 既存スライドのノート（発表原稿）を退避 — 再生成で消さないため。
# スライドの順序が変わると位置がずれる点に注意。
_saved_notes = {}
if os.path.exists(OUT):
    try:
        _old = Presentation(OUT)
        for _i, _s in enumerate(_old.slides):
            if _s.has_notes_slide:
                _t = _s.notes_slide.notes_text_frame.text
                if _t.strip():
                    _saved_notes[_i] = _t
        if _saved_notes:
            print("既存ノートを %d 件退避しました" % len(_saved_notes))
    except Exception as _e:
        print("ノート退避をスキップ:", _e)

prs = Presentation()
prs.slide_width  = W
prs.slide_height = H
blank = prs.slide_layouts[6]
ANS = 'http://schemas.openxmlformats.org/drawingml/2006/main'

# ── XML helpers ──────────────────────────────────────────────
def _alpha(shape, pct):
    solid = shape._element.find('.//{%s}solidFill' % ANS)
    if solid is None:
        return
    for tag in ('srgbClr', 'sysClr', 'schemeClr'):
        clr = solid.find('{%s}%s' % (ANS, tag))
        if clr is not None:
            for old in clr.findall('{%s}alpha' % ANS):
                clr.remove(old)
            a = etree.SubElement(clr, '{%s}alpha' % ANS)
            a.set('val', str(int(pct * 1000)))
            return

def _line_alpha(shape, pct):
    ln = shape._element.find('.//{%s}ln' % ANS)
    if ln is None:
        return
    solid = ln.find('{%s}solidFill' % ANS)
    if solid is None:
        return
    for tag in ('srgbClr', 'sysClr', 'schemeClr'):
        clr = solid.find('{%s}%s' % (ANS, tag))
        if clr is not None:
            for old in clr.findall('{%s}alpha' % ANS):
                clr.remove(old)
            a = etree.SubElement(clr, '{%s}alpha' % ANS)
            a.set('val', str(int(pct * 1000)))
            return

def round_it(shape, adj=12000):
    """Apple 風の角丸に変換。adj: 0-50000 (大きいほど丸い)"""
    g = shape._element.find('.//{%s}prstGeom' % ANS)
    if g is None:
        return
    g.set('prst', 'roundRect')
    av = g.find('{%s}avLst' % ANS)
    if av is None:
        av = etree.SubElement(g, '{%s}avLst' % ANS)
    for old in list(av):
        av.remove(old)
    gd = etree.SubElement(av, '{%s}gd' % ANS)
    gd.set('name', 'adj')
    gd.set('fmla', 'val %d' % adj)

def grad_fill(shape, c1, c2, angle=90):
    """2色グラデーション塗り。angle: 度 (0=左→右, 90=上→下)"""
    f = shape.fill
    f.gradient()
    try:
        f.gradient_angle = angle
    except Exception:
        pass
    f.gradient_stops[0].position = 0.0
    f.gradient_stops[0].color.rgb = c1
    f.gradient_stops[1].position = 1.0
    f.gradient_stops[1].color.rgb = c2

def to_back(slide, shape):
    sp = shape._element
    sp.getparent().remove(sp)
    slide.shapes._spTree.insert(2, sp)

# ── Slide primitives ─────────────────────────────────────────
def add_slide(grad=True):
    s = prs.slides.add_slide(blank)
    s.background.fill.solid()
    s.background.fill.fore_color.rgb = BLACK
    if grad:
        bg = s.shapes.add_shape(1, 0, 0, W, H)
        grad_fill(bg, DBLUE, BLACK, angle=130)
        bg.line.fill.background()
        to_back(s, bg)
    return s

def glass(slide, x, y, w, h, fill=GLASS, fa=55, ba=14, radius=12000):
    sh = slide.shapes.add_shape(1, x, y, w, h)
    sh.fill.solid()
    sh.fill.fore_color.rgb = fill
    _alpha(sh, fa)
    sh.line.color.rgb = WHITE
    sh.line.width = Pt(1)
    _line_alpha(sh, ba)
    round_it(sh, radius)
    sh.shadow.inherit = False
    return sh

def grad_card(slide, x, y, w, h, c1, c2, angle=120, radius=14000, alpha=None):
    sh = slide.shapes.add_shape(1, x, y, w, h)
    grad_fill(sh, c1, c2, angle)
    sh.line.fill.background()
    round_it(sh, radius)
    sh.shadow.inherit = False
    return sh

def bar(slide, x, y, w, h=Inches(0.06), c1=BLUE, c2=LBLUE, radius=50000):
    sh = slide.shapes.add_shape(1, x, y, w, h)
    grad_fill(sh, c1, c2, angle=0)
    sh.line.fill.background()
    round_it(sh, radius)
    sh.shadow.inherit = False
    return sh

def pill(slide, x, y, w, h, label, color):
    sh = slide.shapes.add_shape(1, x, y, w, h)
    sh.fill.solid()
    sh.fill.fore_color.rgb = color
    _alpha(sh, 18)
    sh.line.color.rgb = color
    sh.line.width = Pt(1)
    _line_alpha(sh, 55)
    round_it(sh, 50000)
    sh.shadow.inherit = False
    tf = sh.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    r = p.add_run(); r.text = label
    r.font.name = "Inter"; r.font.size = Pt(15); r.font.bold = True
    r.font.color.rgb = color

def txt(slide, text, x, y, w, h, sz=18, bold=False, col=WHITE,
        align=PP_ALIGN.LEFT, font="Inter", italic=False, anchor=None, spacing=None):
    box = slide.shapes.add_textbox(x, y, w, h)
    tf = box.text_frame
    tf.word_wrap = True
    if anchor:
        tf.vertical_anchor = anchor
    lines = text.split("\n")
    for i, ln in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        if spacing:
            p.line_spacing = spacing
        r = p.add_run(); r.text = ln
        r.font.name = font; r.font.size = Pt(sz); r.font.bold = bold
        r.font.italic = italic; r.font.color.rgb = col
    return box

def mtxt(slide, lines, x, y, w, h, dsz=16, dcol=WHITE, spacing=1.25):
    box = slide.shapes.add_textbox(x, y, w, h)
    tf = box.text_frame
    tf.word_wrap = True
    first = True
    for item in lines:
        if isinstance(item, str):
            t, sz, bd, cl = item, dsz, False, dcol
        else:
            t  = item[0]
            sz = item[1] if len(item) > 1 else dsz
            bd = item[2] if len(item) > 2 else False
            cl = item[3] if len(item) > 3 else dcol
        p = tf.paragraphs[0] if first else tf.add_paragraph()
        first = False
        p.line_spacing = spacing
        r = p.add_run(); r.text = t
        r.font.name = "Inter"; r.font.size = Pt(sz); r.font.bold = bd
        r.font.color.rgb = cl

def kicker(slide, text):
    txt(slide, text, ML, Inches(0.55), CW, Inches(0.34),
        sz=12, bold=True, col=BLUE, font="Inter")

import os
ASSET = os.path.join(os.path.dirname(os.path.abspath(__file__)), "assets")
def pic(slide, name, x, y, w, h):
    path = os.path.join(ASSET, name)
    return slide.shapes.add_picture(path, x, y, width=w, height=h)

# ═══════════════════════════════════════════════════════════════
# 1 — TITLE
# ═══════════════════════════════════════════════════════════════
s = add_slide()
txt(s, "HowTune", ML, Inches(2.35), Inches(11), Inches(1.8),
    sz=96, bold=True, col=WHITE)
bar(s, ML+Inches(0.05), Inches(4.15), Inches(3.6), h=Inches(0.07))
txt(s, "何を聴くかではなく、どう聴いているか。",
    ML, Inches(4.45), Inches(11), Inches(0.7), sz=22, col=LGRAY)
txt(s, "Team Othello   ·   Engineer Guild Hackathon 2026/05",
    ML, Inches(6.55), Inches(11), Inches(0.45), sz=13, col=GRAY)

# ═══════════════════════════════════════════════════════════════
# 2 — CORE IDEA: 曲の特定の瞬間でつながる
# ═══════════════════════════════════════════════════════════════
s = add_slide()
kicker(s, "THE CORE IDEA")
txt(s, "つながるのは、曲の「この一点」。",
    ML, Inches(1.0), CW, Inches(0.95), sz=44, bold=True)
bar(s, ML+Inches(0.05), Inches(2.05), Inches(2.6))

# 楽曲タイムラインのカード
cy, ch = Inches(2.6), Inches(3.0)
glass(s, ML, cy, CW, ch, fa=45)
txt(s, "♪  Blinding Lights — The Weeknd",
    ML+Inches(0.55), cy+Inches(0.3), Inches(8), Inches(0.5),
    sz=18, bold=True, col=WHITE)

# タイムライン（反応地点 = 時刻 + 歌詞）
ts_in, tw_in = 1.55, 10.233          # track start / width (inch)
mfrac = 0.62                          # 1:18 の位置
mx_in = ts_in + tw_in * mfrac
track_y = cy + Inches(1.35)
base = s.shapes.add_shape(1, Inches(ts_in), track_y, Inches(tw_in), Inches(0.14))
base.fill.solid(); base.fill.fore_color.rgb = RGBColor(0x33,0x33,0x38)
base.line.fill.background(); round_it(base, 50000); base.shadow.inherit = False
bar(s, Inches(ts_in), track_y, Inches(tw_in*mfrac), h=Inches(0.14))
# 再生ヘッド（マーカー）
dot = s.shapes.add_shape(9, Inches(mx_in-0.14), track_y-Inches(0.07), Inches(0.28), Inches(0.28))
dot.fill.solid(); dot.fill.fore_color.rgb = LBLUE
dot.line.color.rgb = WHITE; dot.line.width = Pt(1.5); dot.shadow.inherit = False
txt(s, "1:18", Inches(mx_in-0.55), track_y-Inches(0.6), Inches(1.1), Inches(0.4),
    sz=16, bold=True, col=LBLUE, align=PP_ALIGN.CENTER)
# その瞬間の歌詞
txt(s, "“I said, ooh, I'm blinded by the lights”",
    ML+Inches(0.55), cy+Inches(2.0), CW-Inches(1.1), Inches(0.5),
    sz=19, col=WHITE, align=PP_ALIGN.CENTER, italic=True)
txt(s, "●  ●  ●  ●    同じ瞬間に反応した人と出会う",
    ML+Inches(0.55), cy+Inches(2.5), CW-Inches(1.1), Inches(0.4),
    sz=15, bold=True, col=LBLUE, align=PP_ALIGN.CENTER)

txt(s, "曲全体でも、ジャンルでもない。特定の瞬間の歌詞 × 反応で、人とつながる。",
    ML, Inches(5.85), CW, Inches(0.6), sz=18, col=LGRAY)

# ═══════════════════════════════════════════════════════════════
# 3 — WHY NOW (旧→新を2カードで)
# ═══════════════════════════════════════════════════════════════
s = add_slide()
kicker(s, "WHY NOW")
txt(s, "AIが、価値をひっくり返す。",
    ML, Inches(1.05), CW, Inches(1.0), sz=48, bold=True)
bar(s, ML+Inches(0.05), Inches(2.15), Inches(2.6))

gap = Inches(0.5)
cw2 = (CW - gap) / 2
glass(s, ML, Inches(2.85), cw2, Inches(3.6), fa=40)
txt(s, "これまで", ML+Inches(0.5), Inches(3.2), cw2-Inches(1), Inches(0.5),
    sz=15, bold=True, col=GRAY)
txt(s, "何を作るか", ML+Inches(0.5), Inches(3.9), cw2-Inches(1), Inches(0.9),
    sz=40, bold=True, col=LGRAY)
txt(s, "スキルの希少性が価値だった。\nAIがそれをコモディティ化する。",
    ML+Inches(0.5), Inches(5.0), cw2-Inches(1), Inches(1.2), sz=16, col=GRAY, spacing=1.3)

c2x = ML + cw2 + gap
grad_card(s, c2x, Inches(2.85), cw2, Inches(3.6), DBLUE, RGBColor(0x0E,0x2A,0x52), angle=130)
txt(s, "これから", c2x+Inches(0.5), Inches(3.2), cw2-Inches(1), Inches(0.5),
    sz=15, bold=True, col=LBLUE)
txt(s, "どう楽しむか", c2x+Inches(0.5), Inches(3.9), cw2-Inches(1), Inches(0.9),
    sz=40, bold=True, col=WHITE)
txt(s, "体験とコミュニティが最も希少な資源になる。\nHowTune はその中心にいる。",
    c2x+Inches(0.5), Inches(5.0), cw2-Inches(1), Inches(1.2), sz=16, col=LGRAY, spacing=1.3)

# ═══════════════════════════════════════════════════════════════
# 4 — PROBLEM (絵文字3枚)
# ═══════════════════════════════════════════════════════════════
s = add_slide()
kicker(s, "THE PROBLEM")
txt(s, "Labeling is not enough.",
    ML, Inches(1.05), CW, Inches(1.0), sz=52, bold=True)
bar(s, ML+Inches(0.05), Inches(2.2), Inches(2.6))
txt(s, "既存サービスは「曲名」でしか繋げない。本当の熱狂は、聴き方の中にある。",
    ML, Inches(2.5), CW, Inches(0.6), sz=18, col=LGRAY)

cards = [
    ("🎵", "歌詞の一節で\n息が止まる"),
    ("🥁", "ベースで\n体が動く"),
    ("✨", "サビ前の溜めで\n鳥肌が立つ"),
]
gap = Inches(0.4)
cwd = (CW - gap*2) / 3
for i, (emo, t) in enumerate(cards):
    x = ML + i*(cwd+gap)
    glass(s, x, Inches(3.5), cwd, Inches(3.0), fa=45)
    txt(s, emo, x, Inches(3.9), cwd, Inches(0.9), sz=52, align=PP_ALIGN.CENTER)
    txt(s, t, x, Inches(5.0), cwd, Inches(1.2), sz=19, bold=True,
        align=PP_ALIGN.CENTER, spacing=1.2)

# ═══════════════════════════════════════════════════════════════
# 5 — VISION
# ═══════════════════════════════════════════════════════════════
s = add_slide()
kicker(s, "THE SOLUTION")
txt(s, "歌詞を、鏡に。",
    ML, Inches(2.1), CW, Inches(1.4), sz=74, bold=True)
bar(s, ML+Inches(0.05), Inches(3.85), Inches(3.2))
mtxt(s, [
    ("センサーは事実を捉える。", 22, False, LGRAY),
    ("AIは断面を差し出す。", 22, False, LGRAY),
    ("意味は、あなたが見つける。", 22, True, WHITE),
], ML, Inches(4.2), CW, Inches(2.5), spacing=1.35)

# ═══════════════════════════════════════════════════════════════
# 6 — CORE FLOW (5枚・絵文字+1語)
# ═══════════════════════════════════════════════════════════════
s = add_slide()
kicker(s, "HOW IT WORKS")
txt(s, "聴くだけで、体験が言葉になる。",
    ML, Inches(1.05), CW, Inches(0.9), sz=42, bold=True)
bar(s, ML+Inches(0.05), Inches(2.1), Inches(3.0))

steps = [
    ("🎧", "聴く"), ("👆", "タップ"), ("🌊", "Groove"),
    ("💬", "AI 対話"), ("🪪", "How Card"),
]
gap = Inches(0.3)
cwd = (CW - gap*4) / 5
for i, (emo, t) in enumerate(steps):
    x = ML + i*(cwd+gap)
    glass(s, x, Inches(2.85), cwd, Inches(3.3), fa=45)
    txt(s, emo, x, Inches(3.35), cwd, Inches(1.0), sz=46, align=PP_ALIGN.CENTER)
    bar(s, x+cwd/2-Inches(0.35), Inches(4.5), Inches(0.7), h=Inches(0.05))
    txt(s, t, x, Inches(4.75), cwd, Inches(0.6), sz=18, bold=True, align=PP_ALIGN.CENTER)
    txt(s, "0%d" % (i+1), x, Inches(5.45), cwd, Inches(0.4), sz=13, col=BLUE, align=PP_ALIGN.CENTER)

# ═══════════════════════════════════════════════════════════════
# 7 — DEMO
# ═══════════════════════════════════════════════════════════════
s = add_slide()
kicker(s, "LIVE DEMO")
txt(s, "DEMO", ML, Inches(2.0), Inches(8), Inches(2.0), sz=110, bold=True)
bar(s, ML+Inches(0.05), Inches(4.3), Inches(4.0), h=Inches(0.08))
txt(s, "実機で、歌詞タップ → AI対話 → HowCard を体験します。",
    ML, Inches(4.65), CW, Inches(0.6), sz=20, col=LGRAY)
txt(s, "▶ ここに実機スクリーンショット / 画面録画を配置",
    ML, Inches(6.4), CW, Inches(0.5), sz=14, col=GRAY, italic=True)

# ═══════════════════════════════════════════════════════════════
# 8 — AI PHILOSOPHY (対比・短文)
# ═══════════════════════════════════════════════════════════════
s = add_slide()
kicker(s, "AI DESIGN")
txt(s, "断定しない。問いかける。",
    ML, Inches(1.05), CW, Inches(1.0), sz=48, bold=True)
bar(s, ML+Inches(0.05), Inches(2.2), Inches(2.6))

gap = Inches(0.5)
cw2 = (CW - gap) / 2
glass(s, ML, Inches(2.9), cw2, Inches(3.4), fa=35)
txt(s, "✕", ML+Inches(0.5), Inches(3.2), Inches(1), Inches(0.7), sz=30, bold=True, col=RED)
txt(s, "断定する AI", ML+Inches(1.3), Inches(3.28), cw2-Inches(1.8), Inches(0.6),
    sz=20, bold=True, col=RED)
mtxt(s, [
    ("「あなたは感動しました」", 18, False, GRAY),
    ("「ベースが好きですね」", 18, False, GRAY),
], ML+Inches(0.5), Inches(4.2), cw2-Inches(1), Inches(1.8), spacing=1.5)

c2x = ML + cw2 + gap
grad_card(s, c2x, Inches(2.9), cw2, Inches(3.4), DBLUE, RGBColor(0x0E,0x2A,0x52), angle=130)
txt(s, "✓", c2x+Inches(0.5), Inches(3.2), Inches(1), Inches(0.7), sz=30, bold=True, col=LBLUE)
txt(s, "HowTune の AI", c2x+Inches(1.3), Inches(3.28), cw2-Inches(1.8), Inches(0.6),
    sz=20, bold=True, col=WHITE)
mtxt(s, [
    ("「ここで反応していましたね」", 18, False, WHITE),
    ("「リズム？それとも歌詞？」", 18, False, WHITE),
], c2x+Inches(0.5), Inches(4.2), cw2-Inches(1), Inches(1.8), spacing=1.5)

# ═══════════════════════════════════════════════════════════════
# 9 — GROOVE (数式 + 6 pill)
# ═══════════════════════════════════════════════════════════════
s = add_slide()
kicker(s, "THE SCIENCE OF GROOVE")
txt(s, "身体反応を、数値に。",
    ML, Inches(1.05), CW, Inches(0.9), sz=46, bold=True)
bar(s, ML+Inches(0.05), Inches(2.1), Inches(2.6))

grad_card(s, ML, Inches(2.6), CW, Inches(1.5), RGBColor(0x06,0x1A,0x36), DBLUE, angle=0)
txt(s, "Groove  =  volume × 0.42  +  pulse × 0.58",
    ML, Inches(2.6), CW, Inches(1.0), sz=28, bold=True, col=LBLUE,
    align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
txt(s, "AirPods 不要。iPhone だけで動く。",
    ML, Inches(3.55), CW, Inches(0.5), sz=14, col=GRAY, align=PP_ALIGN.CENTER)

axes = [("groove",BLUE),("hype",LBLUE),("chill",GREEN),
        ("immersion",PURPLE),("hit",AMBER),("afterglow",PINK)]
gap = Inches(0.25)
pw = (CW - gap*5) / 6
for i,(name,color) in enumerate(axes):
    pill(s, ML+i*(pw+gap), Inches(4.6), pw, Inches(0.7), name, color)
txt(s, "6軸スコアで「聴き方」を多次元に表現する",
    ML, Inches(5.6), CW, Inches(0.5), sz=16, col=LGRAY, align=PP_ALIGN.CENTER)

# ═══════════════════════════════════════════════════════════════
# 10 — FLYWHEEL (横一列・短語)
# ═══════════════════════════════════════════════════════════════
s = add_slide()
kicker(s, "THE MOAT")
txt(s, "使うほど、賢くなる。",
    ML, Inches(1.05), CW, Inches(0.9), sz=48, bold=True)
bar(s, ML+Inches(0.05), Inches(2.15), Inches(2.6))
txt(s, "AI対話の回答が、そのまま学習データになる。",
    ML, Inches(2.45), CW, Inches(0.5), sz=18, col=LGRAY)

nodes = ["タップ","AI 対話","学習","精度向上","共鳴増加"]
gap = Inches(0.55)
nw = (CW - gap*4) / 5
for i, t in enumerate(nodes):
    x = ML + i*(nw+gap)
    grad_card(s, x, Inches(3.5), nw, Inches(2.0),
              RGBColor(0x10,0x10,0x16), RGBColor(0x06,0x16,0x2E), angle=120)
    txt(s, "0%d"%(i+1), x, Inches(3.75), nw, Inches(0.4), sz=13, col=BLUE, align=PP_ALIGN.CENTER)
    txt(s, t, x, Inches(4.3), nw, Inches(0.9), sz=18, bold=True, align=PP_ALIGN.CENTER)
    if i < 4:
        txt(s, "→", x+nw+Inches(0.05), Inches(4.25), gap-Inches(0.1), Inches(0.5),
            sz=22, bold=True, col=BLUE, align=PP_ALIGN.CENTER)
txt(s, "このデータ資産が、競合の参入障壁になる。",
    ML, Inches(6.0), CW, Inches(0.5), sz=16, bold=True, col=LBLUE, align=PP_ALIGN.CENTER)

# ═══════════════════════════════════════════════════════════════
# 11 — TECHNOLOGY (色ドット + 最小テキスト)
# ═══════════════════════════════════════════════════════════════
s = add_slide()
kicker(s, "TECHNOLOGY")
txt(s, "Apple Ecosystem × Claude",
    ML, Inches(1.05), CW, Inches(0.9), sz=44, bold=True)
bar(s, ML+Inches(0.05), Inches(2.1), Inches(3.0))

tech = [
    (BLUE,   "iOS",     "SwiftUI + MusicKit"),
    (LBLUE,  "歌詞",     "Musixmatch API"),
    (GREEN,  "Groove",  "音量 + 本体モーション"),
    (AMBER,  "AI",      "Claude (backend経由)"),
    (PURPLE, "ML",      "Create ML → Core ML"),
    (PINK,   "Backend", "Firestore + Node.js"),
]
gap = Inches(0.4)
cwd = (CW - gap*2) / 3
chd = Inches(1.7)
for i, (color, layer, t) in enumerate(tech):
    ci, ri = i % 3, i // 3
    x = ML + ci*(cwd+gap)
    y = Inches(2.7) + ri*(chd+Inches(0.4))
    glass(s, x, y, cwd, chd, fa=40)
    dot = s.shapes.add_shape(9, x+Inches(0.4), y+Inches(0.4), Inches(0.28), Inches(0.28))
    dot.fill.solid(); dot.fill.fore_color.rgb = color
    dot.line.fill.background(); dot.shadow.inherit = False
    txt(s, layer, x+Inches(0.85), y+Inches(0.32), cwd-Inches(1), Inches(0.5),
        sz=16, bold=True, col=color)
    txt(s, t, x+Inches(0.42), y+Inches(0.95), cwd-Inches(0.7), Inches(0.6),
        sz=16, col=WHITE)

# ═══════════════════════════════════════════════════════════════
# 11.5 — AI-DRIVEN DEVELOPMENT (CodeRabbit 顧客事例風)
# ═══════════════════════════════════════════════════════════════
s = add_slide()
kicker(s, "AI-DRIVEN DEVELOPMENT")
txt(s, "AI を、チームメイトに。",
    ML, Inches(1.0), CW, Inches(0.9), sz=46, bold=True)
bar(s, ML+Inches(0.05), Inches(2.05), Inches(2.6))
txt(s, "設計・実装・レビューの全工程に AI を組み込んだ。",
    ML, Inches(2.35), CW, Inches(0.5), sz=18, col=LGRAY)

# 実績メトリクス 3枚
metrics = [("151", "Commits"), ("41", "Pull Requests"), ("5", "Days")]
gap = Inches(0.4)
mcw = (CW - gap*2) / 3
for i, (num, lab) in enumerate(metrics):
    x = ML + i*(mcw+gap)
    grad_card(s, x, Inches(2.95), mcw, Inches(1.5),
              RGBColor(0x06,0x1A,0x36), DBLUE, angle=120)
    txt(s, num, x, Inches(3.1), mcw, Inches(0.9), sz=48, bold=True, col=LBLUE,
        align=PP_ALIGN.CENTER)
    txt(s, lab, x, Inches(4.0), mcw, Inches(0.4), sz=14, col=LGRAY, align=PP_ALIGN.CENTER)

# ツール 2枚（ロゴ入り）
gap2 = Inches(0.5)
tcw = (CW - gap2) / 2
glass(s, ML, Inches(4.75), tcw, Inches(1.65), fa=45)
pic(s, "anthropic.png", ML+Inches(0.4), Inches(5.05), Inches(1.05), Inches(1.05))
txt(s, "Claude Code", ML+Inches(1.7), Inches(5.05), tcw-Inches(2.0), Inches(0.5),
    sz=20, bold=True, col=WHITE)
txt(s, "設計ドキュメント・実装・スライドを生成", ML+Inches(1.7), Inches(5.55), tcw-Inches(2.0), Inches(0.8),
    sz=14, col=LGRAY, spacing=1.2)

t2x = ML + tcw + gap2
glass(s, t2x, Inches(4.75), tcw, Inches(1.65), fa=45)
pic(s, "coderabbit.png", t2x+Inches(0.4), Inches(5.05), Inches(1.05), Inches(1.05))
txt(s, "CodeRabbit", t2x+Inches(1.7), Inches(5.05), tcw-Inches(2.0), Inches(0.5),
    sz=20, bold=True, col=WHITE)
txt(s, "全 PR を日本語で自動レビュー（8 PR で指摘・改善）", t2x+Inches(1.7), Inches(5.55), tcw-Inches(2.0), Inches(0.8),
    sz=14, col=LGRAY, spacing=1.2)

# CodeRabbit 実績の引用（出典リンク）
txt(s, "CodeRabbit 実績: コードデリバリ 86% 高速化 / レビュー指摘 60% 削減  ―  claude.com/ja/customers/coderabbit",
    ML, Inches(6.65), CW, Inches(0.45), sz=12, col=GRAY, italic=True)

# ═══════════════════════════════════════════════════════════════
# 12 — BUSINESS MODEL (大きく・収益を明示)
# ═══════════════════════════════════════════════════════════════
s = add_slide()
kicker(s, "BUSINESS MODEL")
txt(s, "3つの収益エンジン。",
    ML, Inches(1.05), CW, Inches(0.9), sz=46, bold=True)
bar(s, ML+Inches(0.05), Inches(2.15), Inches(2.6))

biz = [
    (BLUE,   "B2B SaaS",  "アーティスト\nインサイト",
     "身体反応ヒートマップを提供", "月額サブスク"),
    (GREEN,  "Premium",   "セレンディピティ\nマッチング",
     "同じ瞬間に反応した人と出会う", "プレミアム課金"),
    (AMBER,  "P2P Tips",  "チップ循環",
     "発見者を可視化し直接還元", "マージン 10%"),
]
gap = Inches(0.45)
cwd = (CW - gap*2) / 3
for i, (color, tag, title, desc, rev) in enumerate(biz):
    x = ML + i*(cwd+gap)
    glass(s, x, Inches(2.75), cwd, Inches(4.0), fa=45, radius=14000)
    head = grad_card(s, x, Inches(2.75), cwd, Inches(0.9),
                     color, RGBColor(max(color[0]-40,0),max(color[1]-40,0),max(color[2]-40,0)),
                     angle=0, radius=14000)
    txt(s, tag, x, Inches(2.75), cwd, Inches(0.9), sz=20, bold=True, col=WHITE,
        align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    txt(s, title, x+Inches(0.4), Inches(3.95), cwd-Inches(0.8), Inches(1.1),
        sz=22, bold=True, col=WHITE, spacing=1.1)
    txt(s, desc, x+Inches(0.4), Inches(5.05), cwd-Inches(0.8), Inches(0.9),
        sz=15, col=LGRAY, spacing=1.25)
    bar(s, x+Inches(0.4), Inches(5.95), cwd-Inches(0.8), h=Inches(0.04), c1=color, c2=color)
    txt(s, rev, x+Inches(0.4), Inches(6.1), cwd-Inches(0.8), Inches(0.5),
        sz=17, bold=True, col=color)

# ═══════════════════════════════════════════════════════════════
# 13 — TEAM + MVP
# ═══════════════════════════════════════════════════════════════
s = add_slide()
kicker(s, "TEAM & MVP")
txt(s, "実装済み、実機で動く。",
    ML, Inches(1.05), CW, Inches(0.9), sz=48, bold=True)
bar(s, ML+Inches(0.05), Inches(2.15), Inches(2.6))
txt(s, "Team Othello — PM / iOS / Backend / ML・Design",
    ML, Inches(2.5), CW, Inches(0.5), sz=17, col=GRAY)

gap = Inches(0.5)
cw2 = (CW - gap) / 2
grad_card(s, ML, Inches(3.2), cw2, Inches(3.4), RGBColor(0x06,0x24,0x12), RGBColor(0x04,0x12,0x0A), angle=130)
txt(s, "✓  実装済み (P0)", ML+Inches(0.5), Inches(3.45), cw2-Inches(1), Inches(0.5),
    sz=18, bold=True, col=GREEN)
mtxt(s, [
    "曲再生 × MusicKit", "時刻同期歌詞 (Musixmatch)",
    "歌詞タップ → AI対話 (Claude)", "HowCard 生成・保存",
    "Groove レベル表示",
], ML+Inches(0.5), Inches(4.1), cw2-Inches(1), Inches(2.5), dsz=16, spacing=1.4)

c2x = ML + cw2 + gap
glass(s, c2x, Inches(3.2), cw2, Inches(3.4), fa=35)
txt(s, "○  意図的にスコープ外", c2x+Inches(0.5), Inches(3.45), cw2-Inches(1), Inches(0.5),
    sz=18, bold=True, col=GRAY)
mtxt(s, [
    "DM・フォロー・タイムライン", "Spotify 連携",
    "完全な認証フロー", "AirPods 必須のセンサー精度",
    "教師データ収集アプリ",
], c2x+Inches(0.5), Inches(4.1), cw2-Inches(1), Inches(2.5), dsz=16, dcol=GRAY, spacing=1.4)

# ═══════════════════════════════════════════════════════════════
# 14 — CLOSING
# ═══════════════════════════════════════════════════════════════
s = add_slide()
txt(s, "音楽は、\nどう聴くかだ。",
    ML, Inches(2.0), CW, Inches(2.6), sz=78, bold=True, spacing=1.05)
bar(s, ML+Inches(0.05), Inches(4.85), Inches(3.6), h=Inches(0.07))
txt(s, "歌詞が、あなたの How を語り始める。",
    ML, Inches(5.15), CW, Inches(0.6), sz=22, col=LGRAY)
txt(s, "HowTune   ·   Team Othello",
    ML, Inches(6.6), CW, Inches(0.45), sz=13, col=GRAY)

# ═══════════════════════════════════════════════════════════════
# REFERENCES — 参考資料
# ═══════════════════════════════════════════════════════════════
s = add_slide()
kicker(s, "REFERENCES")
txt(s, "参考資料",
    ML, Inches(1.0), CW, Inches(0.9), sz=46, bold=True)
bar(s, ML+Inches(0.05), Inches(2.05), Inches(2.6))

refs = [
    ("coderabbit.png", "CodeRabbit × Claude 顧客事例",
     "claude.com/ja/customers/coderabbit"),
    ("anthropic.png", "Claude Code — AI ペアプログラミング",
     "claude.com/claude-code"),
    ("coderabbit.png", "CodeRabbit — AI コードレビュー",
     "coderabbit.ai"),
]
ry = Inches(2.75)
rh = Inches(1.25)
for i, (logo, title, url) in enumerate(refs):
    y = ry + i*(rh + Inches(0.22))
    glass(s, ML, y, CW, rh, fa=42)
    pic(s, logo, ML+Inches(0.4), y+Inches(0.27), Inches(0.7), Inches(0.7))
    txt(s, title, ML+Inches(1.45), y+Inches(0.25), CW-Inches(2), Inches(0.5),
        sz=20, bold=True, col=WHITE)
    txt(s, url, ML+Inches(1.45), y+Inches(0.72), CW-Inches(2), Inches(0.4),
        sz=15, col=LBLUE)

# ═══════════════════════════════════════════════════════════════
# Q&A APPENDIX
# ═══════════════════════════════════════════════════════════════
def qa(q, answer):
    s = add_slide()
    txt(s, "Q&A", ML, Inches(0.5), Inches(3), Inches(0.36), sz=12, bold=True, col=BLUE)
    glass(s, ML, Inches(0.95), CW, Inches(1.3), fa=45, radius=10000)
    txt(s, q, ML+Inches(0.4), Inches(0.95), CW-Inches(0.8), Inches(1.3),
        sz=22, bold=True, col=WHITE, anchor=MSO_ANCHOR.MIDDLE)
    bar(s, ML+Inches(0.05), Inches(2.5), Inches(3.0))
    mtxt(s, answer, ML, Inches(2.8), CW, Inches(4.0), dsz=17, spacing=1.45)

qa("Q1: AirPods がないと精度が落ちる？", [
    ("A. AirPods はサプリメントに降格済み (ADR-0005)。", 18, True, WHITE),
    "",
    ("デモは iPhone 単体で完全動作。", 17, False, LBLUE),
    "音量 × 本体モーションで Groove を算出。AirPods 接続時は精度が上がるが必須ではない。",
    "ロードマップ: 音量 → AirPods → 個人最適化。",
])
qa("Q2: 音楽・歌詞のライセンスは？", [
    ("A. Apple MusicKit で再生 = ライセンス取得済み。", 18, True, WHITE),
    "",
    "音源は保持・配信せず MusicKit に委任。",
    "歌詞は Musixmatch API（商用時は有償プラン）。",
    "Claude API はバックエンド経由でキーを秘匿。",
])
qa("Q3: ML の精度は現実的？", [
    ("A. MVP はルールベース。精度より体験を優先。", 18, True, WHITE),
    "",
    ("volume × 0.42 + motion × 0.58。", 17, False, LBLUE),
    "誇大な精度は主張しない。",
    "AI対話の回答 → ラベル → Create ML で精度向上。使うほど賢くなる。",
])
qa("Q4: スケールするか / 模倣されないか？", [
    ("A. データモートとネットワーク効果。", 18, True, WHITE),
    "",
    "「歌詞 × 身体反応 × AI対話」のラベル付きデータは、後発が追うのに同量のユーザー時間が必要。",
    ("同じ歌詞に反応した人が増えるほど、マッチング価値が上がる。", 17, False, LBLUE),
])
qa("Q5: 心拍・身体データのプライバシーは？", [
    ("A. プライバシーを設計の中心に。", 18, True, WHITE),
    "",
    "心拍は HealthKit 端末内に留め、Firestore には書き込まない。",
    "LLM キーはバックエンド経由 (ADR-0002)。",
    "最小権限・明示同意・削除手段を非機能要件に定義。",
])
qa("Q6: なぜ今、なぜこのチーム？", [
    ("A. タイミング・技術・チームが揃った。", 18, True, WHITE),
    "",
    ("AIがスキルをコモディティ化した直後 — 「How」のSNSはまだ無い。", 17, False, LBLUE),
    "Apple Ecosystem × Claude が今が最良の組み合わせ。",
    "PM / iOS / Backend / ML の4専門職で、3日で動くものを作った。",
])

# 退避したノートを同じ位置のスライドに復元
_slides = list(prs.slides)
for _i, _t in _saved_notes.items():
    if _i < len(_slides):
        _slides[_i].notes_slide.notes_text_frame.text = _t
if _saved_notes:
    print("ノートを %d 件復元しました" % len(_saved_notes))

prs.save(OUT)
n = len(prs.slides)
print("Saved: %s  (%d slides: 16 main + refs + %d Q&A)" % (OUT, n, n-17))
