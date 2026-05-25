"""
HowTune — How Resonance（共鳴マッチング）デッキ。
本編 (generate_pptx.py) と同じ Apple Keynote 風スタイル。
"""
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from lxml import etree

WHITE  = RGBColor(0xFF, 0xFF, 0xFF)
BLACK  = RGBColor(0x00, 0x00, 0x00)
GRAY   = RGBColor(0x8E, 0x8E, 0x93)
LGRAY  = RGBColor(0xB8, 0xB8, 0xBE)
BLUE   = RGBColor(0x0A, 0x84, 0xFF)
LBLUE  = RGBColor(0x64, 0xD2, 0xFF)
GREEN  = RGBColor(0x30, 0xD1, 0x58)
AMBER  = RGBColor(0xFF, 0x9F, 0x0A)
ORANGE = RGBColor(0xFF, 0x6A, 0x1A)
PURPLE = RGBColor(0xBF, 0x5A, 0xF2)
PINK   = RGBColor(0xFF, 0x37, 0x5F)
GLASS  = RGBColor(0x1C, 0x1C, 0x1E)
DBLUE  = RGBColor(0x06, 0x12, 0x24)
DRED   = RGBColor(0x2A, 0x08, 0x06)

W, H = Inches(13.333), Inches(7.5)
ML = Inches(1.0)
CW = W - ML * 2
OUT = "docs/slides/howtune_resonance.pptx"

prs = Presentation()
prs.slide_width, prs.slide_height = W, H
blank = prs.slide_layouts[6]
ANS = 'http://schemas.openxmlformats.org/drawingml/2006/main'

def _alpha(shape, pct):
    solid = shape._element.find('.//{%s}solidFill' % ANS)
    if solid is None: return
    for tag in ('srgbClr','sysClr','schemeClr'):
        clr = solid.find('{%s}%s' % (ANS, tag))
        if clr is not None:
            for old in clr.findall('{%s}alpha' % ANS): clr.remove(old)
            a = etree.SubElement(clr, '{%s}alpha' % ANS); a.set('val', str(int(pct*1000))); return

def _line_alpha(shape, pct):
    ln = shape._element.find('.//{%s}ln' % ANS)
    if ln is None: return
    solid = ln.find('{%s}solidFill' % ANS)
    if solid is None: return
    for tag in ('srgbClr','sysClr','schemeClr'):
        clr = solid.find('{%s}%s' % (ANS, tag))
        if clr is not None:
            for old in clr.findall('{%s}alpha' % ANS): clr.remove(old)
            a = etree.SubElement(clr, '{%s}alpha' % ANS); a.set('val', str(int(pct*1000))); return

def round_it(shape, adj=12000):
    g = shape._element.find('.//{%s}prstGeom' % ANS)
    if g is None: return
    g.set('prst', 'roundRect')
    av = g.find('{%s}avLst' % ANS)
    if av is None: av = etree.SubElement(g, '{%s}avLst' % ANS)
    for old in list(av): av.remove(old)
    gd = etree.SubElement(av, '{%s}gd' % ANS); gd.set('name','adj'); gd.set('fmla','val %d' % adj)

def grad_fill(shape, c1, c2, angle=90):
    f = shape.fill; f.gradient()
    try: f.gradient_angle = angle
    except Exception: pass
    f.gradient_stops[0].position = 0.0; f.gradient_stops[0].color.rgb = c1
    f.gradient_stops[1].position = 1.0; f.gradient_stops[1].color.rgb = c2

def to_back(slide, shape):
    sp = shape._element; sp.getparent().remove(sp); slide.shapes._spTree.insert(2, sp)

def add_slide(c1=DBLUE, c2=BLACK):
    s = prs.slides.add_slide(blank)
    s.background.fill.solid(); s.background.fill.fore_color.rgb = BLACK
    bg = s.shapes.add_shape(1, 0, 0, W, H); grad_fill(bg, c1, c2, angle=130)
    bg.line.fill.background(); to_back(s, bg); return s

def glass(slide, x, y, w, h, fa=45, radius=14000):
    sh = slide.shapes.add_shape(1, x, y, w, h)
    sh.fill.solid(); sh.fill.fore_color.rgb = GLASS; _alpha(sh, fa)
    sh.line.color.rgb = WHITE; sh.line.width = Pt(1); _line_alpha(sh, 14)
    round_it(sh, radius); sh.shadow.inherit = False; return sh

def grad_card(slide, x, y, w, h, c1, c2, angle=120, radius=16000):
    sh = slide.shapes.add_shape(1, x, y, w, h); grad_fill(sh, c1, c2, angle)
    sh.line.fill.background(); round_it(sh, radius); sh.shadow.inherit = False; return sh

def bar(slide, x, y, w, h=Inches(0.07), c1=ORANGE, c2=PINK):
    sh = slide.shapes.add_shape(1, x, y, w, h); grad_fill(sh, c1, c2, angle=0)
    sh.line.fill.background(); round_it(sh, 50000); sh.shadow.inherit = False; return sh

def dot(slide, x, y, color, d=Inches(0.26)):
    sh = slide.shapes.add_shape(9, x, y, d, d); sh.fill.solid(); sh.fill.fore_color.rgb = color
    sh.line.fill.background(); sh.shadow.inherit = False; return sh

def txt(slide, text, x, y, w, h, sz=18, bold=False, col=WHITE, align=PP_ALIGN.LEFT,
        italic=False, anchor=None, spacing=None):
    box = slide.shapes.add_textbox(x, y, w, h); tf = box.text_frame; tf.word_wrap = True
    if anchor: tf.vertical_anchor = anchor
    for i, ln in enumerate(text.split("\n")):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        if spacing: p.line_spacing = spacing
        r = p.add_run(); r.text = ln
        r.font.name = "Inter"; r.font.size = Pt(sz); r.font.bold = bold
        r.font.italic = italic; r.font.color.rgb = col
    return box

def mtxt(slide, lines, x, y, w, h, dsz=15, dcol=WHITE, spacing=1.4):
    box = slide.shapes.add_textbox(x, y, w, h); tf = box.text_frame; tf.word_wrap = True
    first = True
    for item in lines:
        t, sz, bd, cl = (item, dsz, False, dcol) if isinstance(item, str) else (
            item[0], item[1] if len(item)>1 else dsz, item[2] if len(item)>2 else False,
            item[3] if len(item)>3 else dcol)
        p = tf.paragraphs[0] if first else tf.add_paragraph(); first = False
        p.line_spacing = spacing
        r = p.add_run(); r.text = t
        r.font.name = "Inter"; r.font.size = Pt(sz); r.font.bold = bd; r.font.color.rgb = cl

def kicker(slide, text, col=ORANGE):
    txt(slide, text, ML, Inches(0.55), CW, Inches(0.34), sz=12, bold=True, col=col)

# 1 — TITLE
s = add_slide(DRED, BLACK)
txt(s, "How Resonance", ML, Inches(2.2), Inches(11.3), Inches(1.2), sz=64, bold=True)
bar(s, ML+Inches(0.05), Inches(3.6), Inches(3.6))
txt(s, "同じ瞬間に、火がつく。", ML, Inches(3.9), CW, Inches(0.8), sz=26, col=LGRAY)
txt(s, "曲の「この一点」で反応した人と、リアルタイムにつながる。", ML, Inches(4.7), CW, Inches(0.5), sz=16, col=GRAY)
txt(s, "HowTune  ·  Team Othello", ML, Inches(6.55), CW, Inches(0.45), sz=13, col=GRAY)

# 2 — THE MOMENT
s = add_slide()
kicker(s, "THE MOMENT")
txt(s, "一番動いた瞬間を、捉える。", ML, Inches(1.0), CW, Inches(0.9), sz=44, bold=True)
bar(s, ML+Inches(0.05), Inches(2.1), Inches(2.6))
txt(s, "AirPods の頭部モーションから、再生中に最も大きく動いた瞬間を記録。MLは回さない。", ML, Inches(2.4), CW, Inches(0.6), sz=17, col=LGRAY)
grad_card(s, ML, Inches(3.2), CW, Inches(1.5), RGBColor(0x06,0x1A,0x36), DBLUE, angle=0)
txt(s, "peak = max(interactionIntensity) → その playbackTime", ML, Inches(3.2), CW, Inches(1.5),
    sz=24, bold=True, col=LBLUE, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
txt(s, "1回目は決めうちで「ここ、どうでしたか？」。スケール時に分類モデルへ寄せる前提。",
    ML, Inches(5.0), CW, Inches(0.6), sz=15, col=GRAY, align=PP_ALIGN.CENTER)

# 3 — ASK (2-stage)
s = add_slide()
kicker(s, "ASK → DEEPEN")
txt(s, "決めうち → AIで深掘り。", ML, Inches(1.0), CW, Inches(0.9), sz=44, bold=True)
bar(s, ML+Inches(0.05), Inches(2.1), Inches(2.6))
steps = [("01", "ピーク地点を提示", "「M:SS で一番動いていました」固定文", BLUE),
         ("02", "AIが深掘り", "Claude が文脈を踏まえ詳細に問う", PURPLE),
         ("03", "Howカード化", "コメント + タグで投稿", ORANGE)]
gap = Inches(0.4); cwd = (CW - gap*2)/3
for i,(n,t,d,c) in enumerate(steps):
    x = ML + i*(cwd+gap)
    glass(s, x, Inches(2.9), cwd, Inches(2.4))
    txt(s, n, x, Inches(3.15), cwd, Inches(0.4), sz=13, col=c, align=PP_ALIGN.CENTER)
    txt(s, t, x, Inches(3.6), cwd, Inches(0.6), sz=19, bold=True, align=PP_ALIGN.CENTER)
    bar(s, x+cwd/2-Inches(0.35), Inches(4.3), Inches(0.7), h=Inches(0.05), c1=c, c2=c)
    txt(s, d, x+Inches(0.2), Inches(4.5), cwd-Inches(0.4), Inches(0.8), sz=13, col=LGRAY, align=PP_ALIGN.CENTER, spacing=1.2)

# 4 — RESONANCE (showpiece)
s = add_slide(DRED, BLACK)
kicker(s, "RESONANCE", col=ORANGE)
txt(s, "同じ瞬間に、誰かが現れる。", ML, Inches(1.05), CW, Inches(0.9), sz=46, bold=True)
bar(s, ML+Inches(0.05), Inches(2.2), Inches(3.0))
mtxt(s, [
    ("分子がふわっと現れ、摩擦で熱を帯び、発火する——", 20, False, LGRAY),
    ("",),
    ("🔥 同じ瞬間に反応した人がリアルタイムに出現", 20, True, ORANGE),
    ("✦ 別の場所で反応した人も見える（違う解釈との出会い）", 18, False, LBLUE),
], ML, Inches(2.7), CW, Inches(2.4), spacing=1.5)
txt(s, "Firestore リアルタイム購読（ADR-0006）。how-cards を song_id で購読し ±2.5秒で同地点判定。",
    ML, Inches(5.6), CW, Inches(0.6), sz=14, col=GRAY)

# 5 — CONNECT (DM)
s = add_slide()
kicker(s, "CONNECT")
txt(s, "熱量が、つながる。", ML, Inches(1.0), CW, Inches(0.9), sz=46, bold=True)
bar(s, ML+Inches(0.05), Inches(2.1), Inches(2.6))
gap = Inches(0.5); cw2 = (CW-gap)/2
grad_card(s, ML, Inches(2.7), cw2, Inches(3.4), RGBColor(0x2A,0x10,0x06), DRED, angle=130)
txt(s, "🔥 同地点マッチ", ML+Inches(0.5), Inches(3.0), cw2-Inches(1), Inches(0.5), sz=20, bold=True, col=ORANGE)
mtxt(s, ["同じ瞬間に火がついた相手と", "そのままリアルタイム DM", "", "通常画面でも🔥マークで再会"],
     ML+Inches(0.5), Inches(3.7), cw2-Inches(1), Inches(2.2), dsz=16, spacing=1.5)
c2x = ML+cw2+gap
glass(s, c2x, Inches(2.7), cw2, Inches(3.4))
txt(s, "速さへのこだわり", c2x+Inches(0.5), Inches(3.0), cw2-Inches(1), Inches(0.5), sz=20, bold=True, col=LBLUE)
mtxt(s, ["楽観的更新（送信即反映）", "Firestore スナップショット購読", "参加者のみ read/write（rules）"],
     c2x+Inches(0.5), Inches(3.7), cw2-Inches(1), Inches(2.2), dsz=16, spacing=1.5)

# 6 — TECH / 非破壊
s = add_slide()
kicker(s, "ENGINEERING")
txt(s, "壊さず、足す。", ML, Inches(1.0), CW, Inches(0.9), sz=44, bold=True)
bar(s, ML+Inches(0.05), Inches(2.1), Inches(2.6))
tech = [(BLUE, "Peak Motion", "interactionIntensity ピーク（ML不使用）"),
        (PURPLE, "LLM 深掘り", "既存 Claude 経路を2回目に活用"),
        (ORANGE, "Realtime", "Firestore 購読 + 楽観的DM（ADR-0006）"),
        (GREEN, "非破壊統合", "optional 引数で既存挙動を不変に")]
gap = Inches(0.4); cwd = (CW-gap)/2; chd = Inches(1.5)
for i,(c,t,d) in enumerate(tech):
    ci, ri = i%2, i//2
    x = ML + ci*(cwd+gap); y = Inches(2.8) + ri*(chd+Inches(0.35))
    glass(s, x, y, cwd, chd)
    dot(s, x+Inches(0.4), y+Inches(0.38), c)
    txt(s, t, x+Inches(0.85), y+Inches(0.3), cwd-Inches(1), Inches(0.5), sz=17, bold=True, col=c)
    txt(s, d, x+Inches(0.42), y+Inches(0.85), cwd-Inches(0.7), Inches(0.5), sz=14, col=LGRAY)

# 7 — CLOSING
s = add_slide(DRED, BLACK)
txt(s, "同じ瞬間に、\n火がつく。", ML, Inches(2.1), CW, Inches(2.2), sz=58, bold=True, spacing=1.08)
bar(s, ML+Inches(0.05), Inches(4.7), Inches(3.6))
txt(s, "What ではなく How。曲のこの一点で、人とつながる。", ML, Inches(5.0), CW, Inches(0.6), sz=20, col=LGRAY)
txt(s, "HowTune  ·  Team Othello", ML, Inches(6.55), CW, Inches(0.45), sz=13, col=GRAY)

prs.save(OUT)
print("Saved: %s (%d slides)" % (OUT, len(prs.slides)))
