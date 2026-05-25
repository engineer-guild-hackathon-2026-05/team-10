"""
HowTune — Spec-Driven Development アピール（事実ベース）
本編デッキ (generate_pptx.py) と同じ Apple Keynote 風スタイル。
全ての数値・主張はリポジトリで検証可能なものだけを使う。
"""
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from lxml import etree

BLACK  = RGBColor(0x00, 0x00, 0x00)
WHITE  = RGBColor(0xFF, 0xFF, 0xFF)
GRAY   = RGBColor(0x8E, 0x8E, 0x93)
LGRAY  = RGBColor(0xB8, 0xB8, 0xBE)
BLUE   = RGBColor(0x0A, 0x84, 0xFF)
LBLUE  = RGBColor(0x64, 0xD2, 0xFF)
GREEN  = RGBColor(0x30, 0xD1, 0x58)
RED    = RGBColor(0xFF, 0x45, 0x3A)
AMBER  = RGBColor(0xFF, 0x9F, 0x0A)
PURPLE = RGBColor(0xBF, 0x5A, 0xF2)
PINK   = RGBColor(0xFF, 0x37, 0x5F)
GLASS  = RGBColor(0x1C, 0x1C, 0x1E)
DBLUE  = RGBColor(0x06, 0x12, 0x24)

W  = Inches(13.333)
H  = Inches(7.5)
ML = Inches(1.0)
CW = W - ML * 2

OUT = "docs/slides/howtune_sdd.pptx"

prs = Presentation()
prs.slide_width  = W
prs.slide_height = H
blank = prs.slide_layouts[6]
ANS = 'http://schemas.openxmlformats.org/drawingml/2006/main'

def _alpha(shape, pct):
    solid = shape._element.find('.//{%s}solidFill' % ANS)
    if solid is None:
        return
    for tag in ('srgbClr', 'sysClr', 'schemeClr'):
        clr = solid.find('{%s}%s' % (ANS, tag))
        if clr is not None:
            for old in clr.findall('{%s}alpha' % ANS):
                clr.remove(old)
            a = etree.SubElement(clr, '{%s}alpha' % ANS)
            a.set('val', str(int(pct * 1000)))
            return

def _line_alpha(shape, pct):
    ln = shape._element.find('.//{%s}ln' % ANS)
    if ln is None:
        return
    solid = ln.find('{%s}solidFill' % ANS)
    if solid is None:
        return
    for tag in ('srgbClr', 'sysClr', 'schemeClr'):
        clr = solid.find('{%s}%s' % (ANS, tag))
        if clr is not None:
            for old in clr.findall('{%s}alpha' % ANS):
                clr.remove(old)
            a = etree.SubElement(clr, '{%s}alpha' % ANS)
            a.set('val', str(int(pct * 1000)))
            return

def round_it(shape, adj=12000):
    g = shape._element.find('.//{%s}prstGeom' % ANS)
    if g is None:
        return
    g.set('prst', 'roundRect')
    av = g.find('{%s}avLst' % ANS)
    if av is None:
        av = etree.SubElement(g, '{%s}avLst' % ANS)
    for old in list(av):
        av.remove(old)
    gd = etree.SubElement(av, '{%s}gd' % ANS)
    gd.set('name', 'adj')
    gd.set('fmla', 'val %d' % adj)

def grad_fill(shape, c1, c2, angle=90):
    f = shape.fill
    f.gradient()
    try:
        f.gradient_angle = angle
    except Exception:
        pass
    f.gradient_stops[0].position = 0.0
    f.gradient_stops[0].color.rgb = c1
    f.gradient_stops[1].position = 1.0
    f.gradient_stops[1].color.rgb = c2

def to_back(slide, shape):
    sp = shape._element
    sp.getparent().remove(sp)
    slide.shapes._spTree.insert(2, sp)

def add_slide(grad=True):
    s = prs.slides.add_slide(blank)
    s.background.fill.solid()
    s.background.fill.fore_color.rgb = BLACK
    if grad:
        bg = s.shapes.add_shape(1, 0, 0, W, H)
        grad_fill(bg, DBLUE, BLACK, angle=130)
        bg.line.fill.background()
        to_back(s, bg)
    return s

def glass(slide, x, y, w, h, fill=GLASS, fa=55, ba=14, radius=12000):
    sh = slide.shapes.add_shape(1, x, y, w, h)
    sh.fill.solid()
    sh.fill.fore_color.rgb = fill
    _alpha(sh, fa)
    sh.line.color.rgb = WHITE
    sh.line.width = Pt(1)
    _line_alpha(sh, ba)
    round_it(sh, radius)
    sh.shadow.inherit = False
    return sh

def grad_card(slide, x, y, w, h, c1, c2, angle=120, radius=14000):
    sh = slide.shapes.add_shape(1, x, y, w, h)
    grad_fill(sh, c1, c2, angle)
    sh.line.fill.background()
    round_it(sh, radius)
    sh.shadow.inherit = False
    return sh

def bar(slide, x, y, w, h=Inches(0.06), c1=BLUE, c2=LBLUE, radius=50000):
    sh = slide.shapes.add_shape(1, x, y, w, h)
    grad_fill(sh, c1, c2, angle=0)
    sh.line.fill.background()
    round_it(sh, radius)
    sh.shadow.inherit = False
    return sh

def txt(slide, text, x, y, w, h, sz=18, bold=False, col=WHITE,
        align=PP_ALIGN.LEFT, font="Inter", italic=False, anchor=None, spacing=None):
    box = slide.shapes.add_textbox(x, y, w, h)
    tf = box.text_frame
    tf.word_wrap = True
    if anchor:
        tf.vertical_anchor = anchor
    lines = text.split("\n")
    for i, ln in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        if spacing:
            p.line_spacing = spacing
        r = p.add_run(); r.text = ln
        r.font.name = font; r.font.size = Pt(sz); r.font.bold = bold
        r.font.italic = italic; r.font.color.rgb = col
    return box

def mtxt(slide, lines, x, y, w, h, dsz=16, dcol=WHITE, spacing=1.25):
    box = slide.shapes.add_textbox(x, y, w, h)
    tf = box.text_frame
    tf.word_wrap = True
    first = True
    for item in lines:
        if isinstance(item, str):
            t, sz, bd, cl = item, dsz, False, dcol
        else:
            t  = item[0]
            sz = item[1] if len(item) > 1 else dsz
            bd = item[2] if len(item) > 2 else False
            cl = item[3] if len(item) > 3 else dcol
        p = tf.paragraphs[0] if first else tf.add_paragraph()
        first = False
        p.line_spacing = spacing
        r = p.add_run(); r.text = t
        r.font.name = "Inter"; r.font.size = Pt(sz); r.font.bold = bd
        r.font.color.rgb = cl

def kicker(slide, text):
    txt(slide, text, ML, Inches(0.55), CW, Inches(0.34),
        sz=12, bold=True, col=BLUE, font="Inter")

def dot(slide, x, y, color, d=Inches(0.26)):
    sh = slide.shapes.add_shape(9, x, y, d, d)
    sh.fill.solid(); sh.fill.fore_color.rgb = color
    sh.line.fill.background(); sh.shadow.inherit = False
    return sh

# 1 — TITLE
s = add_slide()
txt(s, "Spec-Driven Development", ML, Inches(2.3), Inches(11.3), Inches(1.3),
    sz=60, bold=True, col=WHITE)
bar(s, ML+Inches(0.05), Inches(3.7), Inches(3.6), h=Inches(0.07))
txt(s, "事実で語る、3日間の開発プロセス。", ML, Inches(4.0), CW, Inches(0.7),
    sz=24, col=LGRAY)
txt(s, "盛らない。全てリポジトリで検証可能な証跡。", ML, Inches(4.7), CW, Inches(0.5),
    sz=16, col=GRAY)
txt(s, "HowTune  ·  Team Othello", ML, Inches(6.55), CW, Inches(0.45),
    sz=13, col=GRAY)

# 2 — ARTIFACT TRAIL
s = add_slide()
kicker(s, "THE EVIDENCE")
txt(s, "証跡は、全部ある。", ML, Inches(1.0), CW, Inches(0.95), sz=46, bold=True)
bar(s, ML+Inches(0.05), Inches(2.05), Inches(2.6))
metrics = [("6", "機能スペック", "requirements→design→tasklist", BLUE),
           ("5", "ADR", "不可逆な意思決定の記録", PURPLE),
           ("132", "タスク完了管理", "未完了ゼロ・振り返り記述済み", GREEN),
           ("12", "永続ドキュメント", "実装に同期", AMBER)]
gap = Inches(0.4)
cw = (CW - gap*3) / 4
for i, (num, lab, desc, col) in enumerate(metrics):
    x = ML + i*(cw+gap)
    glass(s, x, Inches(2.7), cw, Inches(2.7), fa=45)
    txt(s, num, x, Inches(2.95), cw, Inches(1.0), sz=52, bold=True, col=col,
        align=PP_ALIGN.CENTER)
    txt(s, lab, x, Inches(4.0), cw, Inches(0.5), sz=16, bold=True, col=WHITE,
        align=PP_ALIGN.CENTER)
    txt(s, desc, x+Inches(0.15), Inches(4.55), cw-Inches(0.3), Inches(0.7), sz=12,
        col=LGRAY, align=PP_ALIGN.CENTER, spacing=1.15)
txt(s, "すべて .steering/ と docs/ にコミット済み。git log で時系列まで追える。",
    ML, Inches(5.75), CW, Inches(0.5), sz=15, col=LGRAY, align=PP_ALIGN.CENTER)

# 3 — SPEC STRUCTURE
s = add_slide()
kicker(s, "HOW WE WORKED")
txt(s, "1機能 = 1スペックの流れ。", ML, Inches(1.0), CW, Inches(0.9), sz=44, bold=True)
bar(s, ML+Inches(0.05), Inches(2.05), Inches(3.0))
steps = [("requirements", "何を・なぜ", BLUE),
         ("design", "どう作るか", LBLUE),
         ("tasklist", "分解して管理", GREEN),
         ("振り返り", "計画との差分を記録", AMBER)]
gap = Inches(0.45)
cwd = (CW - gap*3) / 4
for i, (name, desc, col) in enumerate(steps):
    x = ML + i*(cwd+gap)
    grad_card(s, x, Inches(2.8), cwd, Inches(2.4),
              RGBColor(0x10,0x10,0x16), RGBColor(0x06,0x16,0x2E), angle=120)
    txt(s, "0%d" % (i+1), x, Inches(3.05), cwd, Inches(0.4), sz=13, col=col,
        align=PP_ALIGN.CENTER)
    txt(s, name, x, Inches(3.5), cwd, Inches(0.6), sz=19, bold=True, col=WHITE,
        align=PP_ALIGN.CENTER)
    bar(s, x+cwd/2-Inches(0.35), Inches(4.25), Inches(0.7), h=Inches(0.05), c1=col, c2=col)
    txt(s, desc, x, Inches(4.45), cwd, Inches(0.6), sz=13, col=LGRAY,
        align=PP_ALIGN.CENTER)
    if i < 3:
        txt(s, "→", x+cwd+Inches(0.02), Inches(3.75), gap-Inches(0.04), Inches(0.5),
            sz=22, bold=True, col=BLUE, align=PP_ALIGN.CENTER)
txt(s, "振り返りの「計画と異なった点」まで全スペックに実記述 — 逸脱を隠さず残す。",
    ML, Inches(5.6), CW, Inches(0.5), sz=15, col=LGRAY, align=PP_ALIGN.CENTER)

# 4 — FEATURE x EVIDENCE MAP
s = add_slide()
kicker(s, "COVERAGE — HONEST")
txt(s, "コア体験に、SDDを集中。", ML, Inches(1.0), CW, Inches(0.9), sz=44, bold=True)
bar(s, ML+Inches(0.05), Inches(2.05), Inches(2.6))
gap = Inches(0.5)
cw2 = (CW - gap) / 2
grad_card(s, ML, Inches(2.7), cw2, Inches(3.7), RGBColor(0x06,0x24,0x12), RGBColor(0x04,0x12,0x0A), angle=130)
txt(s, "✓  証跡あり", ML+Inches(0.5), Inches(2.95), cw2-Inches(1), Inches(0.5),
    sz=18, bold=True, col=GREEN)
mtxt(s, [
    "歌詞表示UI  (lyrics-ui)",
    "AirPods反応検出  (airpods-interaction)",
    "AI深掘り対話  (how-chat-deepening)",
    "コメント/認証  (firestore-comments-auth)",
    "収益モデル  (revenue-model)",
    "コード品質監査  (code-quality-audit)",
], ML+Inches(0.5), Inches(3.6), cw2-Inches(1), Inches(2.7), dsz=15, spacing=1.45)
c2x = ML + cw2 + gap
glass(s, c2x, Inches(2.7), cw2, Inches(3.7), fa=35)
txt(s, "○  スペックなし（意図的）", c2x+Inches(0.5), Inches(2.95), cw2-Inches(1), Inches(0.5),
    sz=18, bold=True, col=GRAY)
mtxt(s, [
    "MusicFeed / ForYou",
    "NowPlaying / ClipCreation",
    "Onboarding",
    "",
    ("→ フィード系UIはMVP優先で", 14, False, LGRAY),
    ("　コア体験にリソースを集中した判断", 14, False, LGRAY),
], c2x+Inches(0.5), Inches(3.6), cw2-Inches(1), Inches(2.7), dsz=15, dcol=GRAY, spacing=1.45)
txt(s, "「全機能」ではなく「コア体験の全機能」— ここが事実。",
    ML, Inches(6.6), CW, Inches(0.45), sz=15, bold=True, col=LBLUE, align=PP_ALIGN.CENTER)

# 5 — SPEC <-> CODE FIDELITY
s = add_slide()
kicker(s, "SPEC <-> CODE")
txt(s, "仕様は、実装と紐づく。", ML, Inches(1.0), CW, Inches(0.9), sz=44, bold=True)
bar(s, ML+Inches(0.05), Inches(2.05), Inches(2.6))
grad_card(s, ML, Inches(2.6), CW, Inches(1.3), RGBColor(0x06,0x1A,0x36), DBLUE, angle=0)
txt(s, "11 / 11", ML+Inches(0.5), Inches(2.75), Inches(3), Inches(1.0),
    sz=40, bold=True, col=LBLUE, anchor=MSO_ANCHOR.MIDDLE)
txt(s, "仕様が名指しした主要コンポーネントは、すべて実コードに存在",
    ML+Inches(3.6), Inches(2.75), CW-Inches(4), Inches(1.0), sz=16, col=WHITE,
    anchor=MSO_ANCHOR.MIDDLE, spacing=1.2)
gap = Inches(0.5)
cw2 = (CW - gap) / 2
glass(s, ML, Inches(4.15), cw2, Inches(2.2), fa=40)
txt(s, "✓ 挙動まで一致（直近機能）", ML+Inches(0.45), Inches(4.4), cw2-Inches(0.9), Inches(0.5),
    sz=16, bold=True, col=GREEN)
mtxt(s, [
    "ターン数 = 2 (maximumDialogueTurns)",
    "ReactionEvent.score フィールド",
    "ChatPayload.dominantAxis",
    "CMHeadphoneMotionManager 実使用",
], ML+Inches(0.45), Inches(5.0), cw2-Inches(0.9), Inches(1.3), dsz=13, spacing=1.35)
c2x = ML + cw2 + gap
glass(s, c2x, Inches(4.15), cw2, Inches(2.2), fa=40)
txt(s, "△ 実装が追い越した（初期UI）", c2x+Inches(0.45), Inches(4.4), cw2-Inches(0.9), Inches(0.5),
    sz=16, bold=True, col=AMBER)
mtxt(s, [
    "初日の歌詞UI仕様の一部設計は、",
    "その後のHomeView作り直しで変更",
    "",
    ("→ 逸脱は振り返りに記録済み", 13, True, LGRAY),
], c2x+Inches(0.45), Inches(5.0), cw2-Inches(0.9), Inches(1.3), dsz=13, dcol=LGRAY, spacing=1.35)

# 6 — ADAPTIVE SDD x AI
s = add_slide()
kicker(s, "WHY IT WORKS")
txt(s, "適応型SDD × AI。", ML, Inches(1.0), CW, Inches(0.9), sz=44, bold=True)
bar(s, ML+Inches(0.05), Inches(2.05), Inches(2.6))
txt(s, "3日に最適化した軽量プロセス。重い仕様書を先に固めるのではなく——",
    ML, Inches(2.35), CW, Inches(0.5), sz=17, col=LGRAY)
cards = [
    (BLUE,   "Claude Code", "仕様・設計・tasklist を生成"),
    (GREEN,  "steering スキル", "tasklist 更新を強制し進捗を可視化"),
    (AMBER,  "CodeRabbit", "全PRを自動レビュー"),
]
gap = Inches(0.4)
cwd = (CW - gap*2) / 3
for i, (col, t, d) in enumerate(cards):
    x = ML + i*(cwd+gap)
    glass(s, x, Inches(3.05), cwd, Inches(2.0), fa=45)
    dot(s, x+Inches(0.4), Inches(3.35), col)
    txt(s, t, x+Inches(0.82), Inches(3.3), cwd-Inches(1), Inches(0.5), sz=17, bold=True, col=col)
    txt(s, d, x+Inches(0.42), Inches(3.95), cwd-Inches(0.7), Inches(0.9), sz=14, col=LGRAY, spacing=1.2)
txt(s, "「軽量スペックで認識合わせ → 実装 → 逸脱を記録」。不可逆な判断は ADR で理由ごと残す。",
    ML, Inches(5.4), CW, Inches(0.6), sz=16, col=LGRAY, align=PP_ALIGN.CENTER)
txt(s, "= SDD と AI活用は、同じ1本のストーリー。",
    ML, Inches(6.0), CW, Inches(0.5), sz=17, bold=True, col=LBLUE, align=PP_ALIGN.CENTER)

# 7 — CLOSING
s = add_slide()
txt(s, "証跡は全部ある。\n逸脱は正直に記録。", ML, Inches(2.0), CW, Inches(2.2),
    sz=56, bold=True, spacing=1.08)
bar(s, ML+Inches(0.05), Inches(4.6), Inches(3.6), h=Inches(0.07))
txt(s, "重要な判断は ADR に、理由ごと。— 深掘りされても崩れない。",
    ML, Inches(4.9), CW, Inches(0.6), sz=20, col=LGRAY)
txt(s, "HowTune  ·  Team Othello", ML, Inches(6.55), CW, Inches(0.45), sz=13, col=GRAY)

prs.save(OUT)
print("Saved: %s  (%d slides)" % (OUT, len(prs.slides)))
