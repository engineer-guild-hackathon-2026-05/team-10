"""Shared PowerPoint drawing helpers for HowTune slide generators."""
import logging

from lxml import etree
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt

ANS = "http://schemas.openxmlformats.org/drawingml/2006/main"
logger = logging.getLogger(__name__)

_ctx = {}


def configure(
    *,
    presentation,
    blank_layout,
    width,
    height,
    margin_left,
    content_width,
    white,
    black,
    glass_fill,
    default_gradient_start,
    default_gradient_end,
    default_glass_alpha,
    default_glass_radius,
    default_bar_height,
    default_bar_start,
    default_bar_end,
    default_kicker_color,
    default_glass_border_alpha=14,
):
    _ctx.update(
        presentation=presentation,
        blank_layout=blank_layout,
        width=width,
        height=height,
        margin_left=margin_left,
        content_width=content_width,
        white=white,
        black=black,
        glass_fill=glass_fill,
        default_gradient_start=default_gradient_start,
        default_gradient_end=default_gradient_end,
        default_glass_alpha=default_glass_alpha,
        default_glass_radius=default_glass_radius,
        default_glass_border_alpha=default_glass_border_alpha,
        default_bar_height=default_bar_height,
        default_bar_start=default_bar_start,
        default_bar_end=default_bar_end,
        default_kicker_color=default_kicker_color,
    )


def _value(name):
    if name not in _ctx:
        raise RuntimeError("pptx_utils.configure() must be called before using helpers")
    return _ctx[name]


def _alpha(shape, pct):
    solid = shape._element.find(".//{%s}solidFill" % ANS)
    if solid is None:
        return
    for tag in ("srgbClr", "sysClr", "schemeClr"):
        clr = solid.find("{%s}%s" % (ANS, tag))
        if clr is not None:
            for old in clr.findall("{%s}alpha" % ANS):
                clr.remove(old)
            a = etree.SubElement(clr, "{%s}alpha" % ANS)
            a.set("val", str(int(pct * 1000)))
            return


def _line_alpha(shape, pct):
    ln = shape._element.find(".//{%s}ln" % ANS)
    if ln is None:
        return
    solid = ln.find("{%s}solidFill" % ANS)
    if solid is None:
        return
    for tag in ("srgbClr", "sysClr", "schemeClr"):
        clr = solid.find("{%s}%s" % (ANS, tag))
        if clr is not None:
            for old in clr.findall("{%s}alpha" % ANS):
                clr.remove(old)
            a = etree.SubElement(clr, "{%s}alpha" % ANS)
            a.set("val", str(int(pct * 1000)))
            return


def round_it(shape, adj=12000):
    g = shape._element.find(".//{%s}prstGeom" % ANS)
    if g is None:
        return
    g.set("prst", "roundRect")
    av = g.find("{%s}avLst" % ANS)
    if av is None:
        av = etree.SubElement(g, "{%s}avLst" % ANS)
    for old in list(av):
        av.remove(old)
    gd = etree.SubElement(av, "{%s}gd" % ANS)
    gd.set("name", "adj")
    gd.set("fmla", "val %d" % adj)


def grad_fill(shape, c1, c2, angle=90):
    fill = shape.fill
    fill.gradient()
    try:
        fill.gradient_angle = angle
    except (AttributeError, TypeError) as exc:
        logger.debug("gradient_angle=%s is not supported: %s", angle, exc)
    fill.gradient_stops[0].position = 0.0
    fill.gradient_stops[0].color.rgb = c1
    fill.gradient_stops[1].position = 1.0
    fill.gradient_stops[1].color.rgb = c2


def to_back(slide, shape):
    sp = shape._element
    sp.getparent().remove(sp)
    slide.shapes._spTree.insert(2, sp)


def add_slide(c1=None, c2=None, grad=True):
    prs = _value("presentation")
    blank = _value("blank_layout")
    s = prs.slides.add_slide(blank)
    s.background.fill.solid()
    s.background.fill.fore_color.rgb = _value("black")
    if grad:
        bg = s.shapes.add_shape(1, 0, 0, _value("width"), _value("height"))
        grad_fill(
            bg,
            c1 or _value("default_gradient_start"),
            c2 or _value("default_gradient_end"),
            angle=130,
        )
        bg.line.fill.background()
        to_back(s, bg)
    return s


def glass(slide, x, y, w, h, fill=None, fa=None, ba=None, radius=None):
    sh = slide.shapes.add_shape(1, x, y, w, h)
    sh.fill.solid()
    sh.fill.fore_color.rgb = fill or _value("glass_fill")
    _alpha(sh, _value("default_glass_alpha") if fa is None else fa)
    sh.line.color.rgb = _value("white")
    sh.line.width = Pt(1)
    _line_alpha(sh, _value("default_glass_border_alpha") if ba is None else ba)
    round_it(sh, _value("default_glass_radius") if radius is None else radius)
    sh.shadow.inherit = False
    return sh


def grad_card(slide, x, y, w, h, c1, c2, angle=120, radius=14000):
    sh = slide.shapes.add_shape(1, x, y, w, h)
    grad_fill(sh, c1, c2, angle)
    sh.line.fill.background()
    round_it(sh, radius)
    sh.shadow.inherit = False
    return sh


def bar(slide, x, y, w, h=None, c1=None, c2=None, radius=50000):
    sh = slide.shapes.add_shape(1, x, y, w, h or _value("default_bar_height"))
    grad_fill(sh, c1 or _value("default_bar_start"), c2 or _value("default_bar_end"), angle=0)
    sh.line.fill.background()
    round_it(sh, radius)
    sh.shadow.inherit = False
    return sh


def dot(slide, x, y, color, d=Inches(0.26)):
    sh = slide.shapes.add_shape(9, x, y, d, d)
    sh.fill.solid()
    sh.fill.fore_color.rgb = color
    sh.line.fill.background()
    sh.shadow.inherit = False
    return sh


def txt(
    slide,
    text,
    x,
    y,
    w,
    h,
    sz=18,
    bold=False,
    col=None,
    align=PP_ALIGN.LEFT,
    font="Inter",
    italic=False,
    anchor=None,
    spacing=None,
):
    box = slide.shapes.add_textbox(x, y, w, h)
    tf = box.text_frame
    tf.word_wrap = True
    if anchor:
        tf.vertical_anchor = anchor
    for i, line in enumerate(text.split("\n")):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        if spacing:
            p.line_spacing = spacing
        r = p.add_run()
        r.text = line
        r.font.name = font
        r.font.size = Pt(sz)
        r.font.bold = bold
        r.font.italic = italic
        r.font.color.rgb = col or _value("white")
    return box


def mtxt(slide, lines, x, y, w, h, dsz=15, dcol=None, spacing=1.4):
    box = slide.shapes.add_textbox(x, y, w, h)
    tf = box.text_frame
    tf.word_wrap = True
    first = True
    for item in lines:
        if isinstance(item, str):
            text, sz, bold, color = item, dsz, False, dcol or _value("white")
        else:
            text = item[0]
            sz = item[1] if len(item) > 1 else dsz
            bold = item[2] if len(item) > 2 else False
            color = item[3] if len(item) > 3 else dcol or _value("white")
        p = tf.paragraphs[0] if first else tf.add_paragraph()
        first = False
        p.line_spacing = spacing
        r = p.add_run()
        r.text = text
        r.font.name = "Inter"
        r.font.size = Pt(sz)
        r.font.bold = bold
        r.font.color.rgb = color
    return box


def kicker(slide, text, col=None):
    txt(
        slide,
        text,
        _value("margin_left"),
        Inches(0.55),
        _value("content_width"),
        Inches(0.34),
        sz=12,
        bold=True,
        col=col or _value("default_kicker_color"),
        font="Inter",
    )
